import io
from typing import Optional

from company_llm_rag.logger import get_logger

logger = get_logger(__name__)


def extract_pdf_text(file_bytes: bytes) -> str:
    """
    PDF 바이너리에서 텍스트를 추출합니다.
    스캔본(이미지 PDF)처럼 텍스트 레이어가 없으면 빈 문자열을 반환합니다.

    Args:
        file_bytes: PDF 파일의 바이너리 데이터

    Returns:
        추출된 텍스트 (페이지별 구분)
    """
    try:
        from pypdf import PdfReader
    except ImportError:
        logger.warning("pypdf가 설치되어 있지 않습니다. PDF 파싱을 건너뜁니다.")
        return ""

    try:
        reader = PdfReader(io.BytesIO(file_bytes))
        pages = []
        for i, page in enumerate(reader.pages, 1):
            text = page.extract_text() or ""
            if text.strip():
                pages.append(f"[Page {i}]\n{text.strip()}")
        return "\n\n".join(pages)
    except Exception as e:
        logger.warning(f"PDF 텍스트 추출 실패: {e}")
        return ""


def extract_pptx_text(file_bytes: bytes) -> str:
    """
    PPTX 바이너리에서 텍스트를 추출합니다.
    텍스트가 없는 슬라이드(이미지만 있는 슬라이드 등)는 건너뜁니다.

    Args:
        file_bytes: PPTX 파일의 바이너리 데이터

    Returns:
        추출된 텍스트 (슬라이드별 구분)
    """
    try:
        from pptx import Presentation
    except ImportError:
        logger.warning("python-pptx가 설치되어 있지 않습니다. PPTX 파싱을 건너뜁니다.")
        return ""

    try:
        prs = Presentation(io.BytesIO(file_bytes))
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            texts = [
                shape.text_frame.text.strip()
                for shape in slide.shapes
                if shape.has_text_frame and shape.text_frame.text.strip()
            ]
            if texts:
                slides.append(f"[Slide {i}]\n" + "\n".join(texts))
        return "\n\n".join(slides)
    except Exception as e:
        logger.warning(f"PPTX 텍스트 추출 실패: {e}")
        return ""


def extract_docx_text(file_bytes: bytes) -> str:
    """
    DOCX 바이너리에서 텍스트를 추출합니다.

    Args:
        file_bytes: DOCX 파일의 바이너리 데이터

    Returns:
        추출된 텍스트 (단락별 구분)
    """
    try:
        from docx import Document
    except ImportError:
        logger.warning("python-docx가 설치되어 있지 않습니다. DOCX 파싱을 건너뜁니다.")
        return ""

    try:
        doc = Document(io.BytesIO(file_bytes))
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        return "\n\n".join(paragraphs)
    except Exception as e:
        logger.warning(f"DOCX 텍스트 추출 실패: {e}")
        return ""


def extract_doc_text(file_bytes: bytes) -> str:
    """
    DOC (구 Word 바이너리) 파일에서 텍스트를 추출합니다.
    python-docx로 best-effort 파싱을 시도합니다.
    진짜 바이너리 .doc 포맷은 파싱 실패 시 빈 문자열을 반환합니다.

    Args:
        file_bytes: DOC 파일의 바이너리 데이터

    Returns:
        추출된 텍스트 또는 빈 문자열
    """
    try:
        from docx import Document
    except ImportError:
        logger.warning("python-docx가 설치되어 있지 않습니다. DOC 파싱을 건너뜁니다.")
        return ""

    try:
        doc = Document(io.BytesIO(file_bytes))
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        return "\n\n".join(paragraphs)
    except Exception as e:
        logger.warning(f"DOC 텍스트 추출 실패 (구 바이너리 포맷 미지원): {e}")
        return ""


def extract_xlsx_text(file_bytes: bytes) -> str:
    """
    XLSX 바이너리에서 텍스트를 추출합니다.
    시트별로 데이터를 추출하며, 빈 행은 건너뜁니다.

    Args:
        file_bytes: XLSX 파일의 바이너리 데이터

    Returns:
        추출된 텍스트 (시트별 구분)
    """
    try:
        import openpyxl
    except ImportError:
        logger.warning("openpyxl이 설치되어 있지 않습니다. XLSX 파싱을 건너뜁니다.")
        return ""

    try:
        wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
        sheets = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
                if cells:
                    rows.append("\t".join(cells))
            if rows:
                sheets.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))
        wb.close()
        return "\n\n".join(sheets)
    except Exception as e:
        logger.warning(f"XLSX 텍스트 추출 실패: {e}")
        return ""


def extract_xls_text(file_bytes: bytes) -> str:
    """
    XLS (구 Excel 바이너리) 파일에서 텍스트를 추출합니다.
    xlrd 1.x (xls 지원 버전)가 필요합니다.

    Args:
        file_bytes: XLS 파일의 바이너리 데이터

    Returns:
        추출된 텍스트 (시트별 구분)
    """
    try:
        import xlrd
    except ImportError:
        logger.warning("xlrd가 설치되어 있지 않습니다. XLS 파싱을 건너뜁니다.")
        return ""

    try:
        wb = xlrd.open_workbook(file_contents=file_bytes)
        sheets = []
        for sheet_name in wb.sheet_names():
            ws = wb.sheet_by_name(sheet_name)
            rows = []
            for row_idx in range(ws.nrows):
                cells = [str(ws.cell_value(row_idx, col)).strip()
                         for col in range(ws.ncols)
                         if str(ws.cell_value(row_idx, col)).strip()]
                if cells:
                    rows.append("\t".join(cells))
            if rows:
                sheets.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))
        return "\n\n".join(sheets)
    except Exception as e:
        logger.warning(f"XLS 텍스트 추출 실패: {e}")
        return ""
